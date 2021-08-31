from wand.image import Image
from pptx import Presentation
from pptx.util import Inches

arr = ['image1.jpg','image2.jpg','image3.jpg','image4.jpg','image5.jpg']
for i in range(5):
     with Image(filename=arr[i]) as background:
        with Image(filename='nike_black.png') as watermark:
            background.watermark(image=watermark, transparency=0.25)
        background.save(filename='result' + str(i) + '.jpg')
prs = Presentation()
for i in range(5):
        img_path = 'result' + str(i) + '.jpg'
        blank_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, Inches(1),  Inches(2), Inches(5), Inches(5))
        
prs.save('test.pptx')